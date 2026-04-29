from __future__ import annotations

from pathlib import Path
from typing import Any, Dict, List, Tuple


def _table_to_markdown(rows: List[List[str]]) -> str:
    """행 리스트를 마크다운 표로 변환. 첫 행을 헤더로 사용. 빈 입력은 ''."""
    if not rows:
        return ""
    col_count = max(len(r) for r in rows)
    if col_count == 0:
        return ""

    def _norm(cells: List[str]) -> List[str]:
        padded = (cells + [""] * col_count)[:col_count]
        return [c.replace("|", "\\|").replace("\n", " ").strip() for c in padded]

    header = _norm(rows[0])
    lines = ["| " + " | ".join(header) + " |"]
    lines.append("|" + "|".join(["---"] * col_count) + "|")
    for row in rows[1:]:
        lines.append("| " + " | ".join(_norm(row)) + " |")
    return "\n".join(lines)


def _read_txt(path: Path) -> str:
    return path.read_text(encoding="utf-8", errors="ignore")


def _read_pdf(path: Path) -> str:
    try:
        from pypdf import PdfReader  # type: ignore
    except Exception as e:
        raise RuntimeError(f"PDF 추출을 위해 pypdf가 필요합니다: {e}")

    reader = PdfReader(str(path))
    out: list[str] = []
    for page in reader.pages:
        try:
            out.append(page.extract_text() or "")
        except Exception:
            out.append("")
    return "\n".join(out).strip()


def _read_docx(path: Path) -> str:
    """DOCX: 본문 단락은 단락 단위로, 표는 마크다운 표로 변환.
    단락·표 사이는 \\n\\n으로 단락 분리하여 청킹 시 경계 보존.
    """
    try:
        import docx  # type: ignore  # python-docx
    except Exception as e:
        raise RuntimeError(f"DOCX 추출을 위해 python-docx가 필요합니다: {e}")

    d = docx.Document(str(path))
    parts: list[str] = []

    for p in d.paragraphs:
        text = (p.text or "").strip()
        if text:
            parts.append(text)

    for idx, table in enumerate(d.tables, start=1):
        rows: List[List[str]] = []
        for row in table.rows:
            cells = [c.text.strip() for c in row.cells]
            if any(cells):
                rows.append(cells)
        if rows:
            parts.append(f"## 표 {idx}")
            parts.append(_table_to_markdown(rows))

    return "\n\n".join(parts).strip()


def _read_xlsx(path: Path) -> str:
    """XLSX: 시트별 마크다운 표로 변환. 첫 행을 헤더로 사용.
    시트 간 \\n\\n으로 분리, 시트명을 헤딩으로 명시.
    """
    try:
        import openpyxl  # type: ignore
    except Exception as e:
        raise RuntimeError(f"XLSX 추출을 위해 openpyxl이 필요합니다: {e}")

    wb = openpyxl.load_workbook(str(path), data_only=True)
    parts: list[str] = []
    for ws in wb.worksheets:
        rows: List[List[str]] = []
        for row in ws.iter_rows(values_only=True):
            vals = ["" if v is None else str(v) for v in row]
            if any(v.strip() for v in vals):
                rows.append(vals)
        if not rows:
            continue
        parts.append(f"## 시트: {ws.title}")
        parts.append(_table_to_markdown(rows))
    return "\n\n".join(parts).strip()


def _read_pptx(file_path: str) -> str:
    """PPTX: 슬라이드 단위 단락 분리 + 제목 헤딩 + 표 마크다운 + 발표자 노트.

    각 슬라이드 출력 형식:
        ## 슬라이드 N
        ### {제목}
        본문 도형 텍스트들...
        | 표 | 마크다운 |
        ### 발표자 노트
        노트 텍스트...

    슬라이드 간 \\n\\n으로 단락 분리하여 청킹 시 슬라이드 경계 보존.
    검색 시 슬라이드 번호·제목으로 위치 추적 가능.
    """
    from pptx import Presentation

    prs = Presentation(file_path)
    slide_blocks: List[str] = []

    for idx, slide in enumerate(prs.slides, start=1):
        parts: List[str] = []

        # 1. 슬라이드 번호 헤딩
        parts.append(f"## 슬라이드 {idx}")

        # 2. 슬라이드 제목 추출
        title_shape = None
        try:
            title_shape = slide.shapes.title
        except Exception:
            title_shape = None
        if title_shape is not None:
            title_text = (title_shape.text or "").strip()
            if title_text:
                parts.append(f"### {title_text}")

        # 3. 본문 도형 + 표 (제목 도형은 중복 회피)
        for shape in slide.shapes:
            if title_shape is not None and shape is title_shape:
                continue
            if getattr(shape, "has_table", False):
                rows: List[List[str]] = []
                for row in shape.table.rows:
                    cells = [(cell.text or "").strip() for cell in row.cells]
                    if any(cells):
                        rows.append(cells)
                if rows:
                    parts.append(_table_to_markdown(rows))
            elif hasattr(shape, "text") and shape.text:
                text = shape.text.strip()
                if text:
                    parts.append(text)

        # 4. 발표자 노트
        try:
            if slide.has_notes_slide:
                notes_text = (slide.notes_slide.notes_text_frame.text or "").strip()
                if notes_text:
                    parts.append(f"### 발표자 노트\n{notes_text}")
        except Exception:
            pass

        slide_blocks.append("\n\n".join(parts))

    return "\n\n".join(slide_blocks).strip()


def extract_text_from_file(path: Path) -> Tuple[str, Dict[str, Any]]:
    suffix = path.suffix.lower()
    meta: Dict[str, Any] = {
        "path": str(path),
        "name": path.name,
        "suffix": suffix,
        "size_bytes": path.stat().st_size,
    }
 
    if suffix in {".txt", ".md", ".csv", ".log"}:
        return _read_txt(path), meta
    if suffix == ".pdf":
        return _read_pdf(path), meta
    if suffix == ".docx":
        return _read_docx(path), meta
    if suffix in {".xlsx", ".xlsm"}:
        return _read_xlsx(path), meta
    if suffix == ".pptx":
        return _read_pptx(path), meta

    raise RuntimeError(f"지원하지 않는 파일 형식: {suffix}")


